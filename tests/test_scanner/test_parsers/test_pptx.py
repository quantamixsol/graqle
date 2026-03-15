"""Tests for graqle.scanner.parsers.pptx — PPTX parser."""

# ── graqle:intelligence ──
# module: tests.test_scanner.test_parsers.test_pptx
# risk: LOW (impact radius: 0 modules)
# dependencies: __future__, pathlib, pytest, pptx
# constraints: none
# ── /graqle:intelligence ──

from __future__ import annotations

from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx", reason="python-pptx not installed")

from graqle.scanner.parsers.pptx import PPTXParser


@pytest.fixture
def parser() -> PPTXParser:
    return PPTXParser()


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches

    p = tmp_path / "test.pptx"
    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Test Author"

    # Slide 1: title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Welcome"
    slide1.placeholders[1].text = "A test presentation"

    # Slide 2: content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Architecture"
    slide2.placeholders[1].text = "The auth_service handles JWT validation."

    # Slide 3: table
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Data"
    table = slide3.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "Service"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Auth"
    table.cell(1, 1).text = "Active"
    table.cell(2, 0).text = "DB"
    table.cell(2, 1).text = "Active"

    prs.save(str(p))
    return p


class TestPPTXParser:
    def test_is_available(self, parser: PPTXParser) -> None:
        assert parser.is_available() is True

    def test_supported_extensions(self, parser: PPTXParser) -> None:
        assert ".pptx" in parser.supported_extensions

    def test_parse_basic(self, parser: PPTXParser, sample_pptx: Path) -> None:
        doc = parser.parse(sample_pptx)
        assert doc.format == "pptx"
        assert doc.path == sample_pptx
        assert len(doc.sections) >= 2

    def test_parse_metadata(self, parser: PPTXParser, sample_pptx: Path) -> None:
        doc = parser.parse(sample_pptx)
        assert doc.metadata.get("slide_count") >= 3
        assert doc.metadata.get("title") == "Test Presentation"

    def test_slides_as_sections(self, parser: PPTXParser, sample_pptx: Path) -> None:
        doc = parser.parse(sample_pptx)
        for section in doc.sections:
            assert section.section_type == "slide"
            assert section.page is not None

    def test_table_extraction(self, parser: PPTXParser, sample_pptx: Path) -> None:
        doc = parser.parse(sample_pptx)
        all_tables = []
        for section in doc.sections:
            all_tables.extend(section.tables)
        assert len(all_tables) >= 1
        assert "Service" in all_tables[0]["headers"]

    def test_full_text_content(self, parser: PPTXParser, sample_pptx: Path) -> None:
        doc = parser.parse(sample_pptx)
        assert "auth_service" in doc.full_text

    def test_file_not_found(self, parser: PPTXParser, tmp_path: Path) -> None:
        with pytest.raises(FileNotFoundError):
            parser.parse(tmp_path / "nonexistent.pptx")
