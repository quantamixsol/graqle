"""PPTX presentation parser — requires ``python-pptx``.

Extracts text from slides, including shapes, tables, and notes.
Each slide becomes a section.
"""

# ── graqle:intelligence ──
# module: graqle.scanner.parsers.pptx
# risk: LOW (impact radius: 1 modules)
# consumers: test_pptx
# dependencies: __future__, pathlib, typing, base
# constraints: none
# ── /graqle:intelligence ──

from __future__ import annotations

from pathlib import Path
from typing import Any

from graqle.scanner.parsers.base import (
    BaseDocParser,
    ParsedDocument,
    ParsedSection,
)


class PPTXParser(BaseDocParser):
    """Parse ``.pptx`` files using ``python-pptx``."""

    def parse(self, path: Path) -> ParsedDocument:
        """Parse a PPTX file into a :class:`ParsedDocument`."""
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401

        path = Path(path)
        if not path.is_file():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        metadata: dict[str, Any] = {"source": str(path)}
        parse_errors: list[str] = []

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ValueError(f"Failed to parse PPTX: {exc}") from exc

        # Metadata
        core = prs.core_properties
        if core.title:
            metadata["title"] = core.title
        if core.author:
            metadata["author"] = core.author
        metadata["slide_count"] = len(prs.slides)

        sections: list[ParsedSection] = []
        full_text_parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            tables: list[dict[str, Any]] = []
            slide_title = ""

            for shape in slide.shapes:
                # Title
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # First text frame is usually the title
                        if not slide_title and shape.shape_type is not None:
                            slide_title = text
                        slide_texts.append(text)

                # Tables
                if shape.has_table:
                    try:
                        tbl = shape.table
                        rows_data = []
                        for row in tbl.rows:
                            rows_data.append([cell.text.strip() for cell in row.cells])
                        if len(rows_data) >= 2:
                            tables.append({
                                "headers": rows_data[0],
                                "rows": rows_data[1:],
                            })
                    except Exception as exc:
                        parse_errors.append(f"Slide {slide_num} table: {exc}")

            # Notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes: {notes}]")

            content = "\n".join(slide_texts).strip()
            if content:
                title = slide_title or f"Slide {slide_num}"
                sections.append(
                    ParsedSection(
                        title=title,
                        content=content,
                        level=1,
                        section_type="slide",
                        page=slide_num,
                        tables=tables,
                        code_blocks=[],
                        links=[],
                    )
                )
                full_text_parts.append(content)

        title = metadata.get("title", path.stem)
        full_text = "\n\n".join(full_text_parts)

        return ParsedDocument(
            path=path,
            title=str(title),
            format="pptx",
            sections=sections,
            full_text=full_text,
            metadata=metadata,
            parse_errors=parse_errors,
        )

    def is_available(self) -> bool:
        try:
            from pptx import Presentation  # noqa: F401
            return True
        except ImportError:
            return False

    def missing_dependency_message(self) -> str:
        return "Install python-pptx: pip install graqle[docs]"

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]
